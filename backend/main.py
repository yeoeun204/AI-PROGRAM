from dotenv import load_dotenv
load_dotenv()

import io
import json
import os
from contextlib import asynccontextmanager

import pdfplumber
from fastapi import FastAPI, UploadFile, File, Form, Depends, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from database import get_db, init_db
from models import Lecture, GraphNode, Quiz, QuizAttempt, NodeScore
from graph_builder import build_graph
from quiz_generator import generate_quizzes
from quiz_evaluator import evaluate_text_answer, evaluate_quantitative_answer
from weight_engine import update_node_weight
from error_analyzer import apply_causal_rules, multi_agent_analysis


@asynccontextmanager
async def lifespan(app: FastAPI):
    await init_db()
    yield

app = FastAPI(title="AI 학습 최적화 플랫폼", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

FRONTEND_DIR = os.path.join(os.path.dirname(__file__), "..", "frontend")
if os.path.isdir(FRONTEND_DIR):
    app.mount("/static", StaticFiles(directory=FRONTEND_DIR), name="static")


# ── 헬퍼 ────────────────────────────────────

def _build_tree(nodes: list[GraphNode]) -> list[dict]:
    node_map = {n.id: {
        "id": n.id,
        "concept": n.concept,
        "description": n.description,
        "importance_score": n.importance_score,
        "parent_id": n.parent_id,
        "level": n.level,
        "children": [],
    } for n in nodes}

    roots = []
    for n in nodes:
        d = node_map[n.id]
        if n.parent_id and n.parent_id in node_map:
            node_map[n.parent_id]["children"].append(d)
        else:
            roots.append(d)
    return roots


def _extract_text_from_pdf(data: bytes) -> str:
    parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
    return "\n".join(parts)


# ── 라우팅 ───────────────────────────────────

@app.get("/")
async def root():
    index_path = os.path.join(FRONTEND_DIR, "index.html")
    if os.path.exists(index_path):
        return FileResponse(index_path)
    return {"message": "AI 학습 최적화 플랫폼"}


@app.post("/api/lectures/upload")
async def upload_lecture(
    title: str = Form(...),
    file: UploadFile = File(None),
    text_content: str = Form(None),
    db: AsyncSession = Depends(get_db),
):
    if file:
        raw = await file.read()
        if file.filename.lower().endswith(".pdf"):
            content = _extract_text_from_pdf(raw)
        elif file.filename.lower().endswith((".pptx", ".ppt")):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            content = "\n".join(parts)
        else:
            content = raw.decode("utf-8", errors="ignore")
    elif text_content:
        content = text_content
    else:
        raise HTTPException(status_code=400, detail="파일 또는 텍스트를 입력하세요.")

    if len(content.strip()) < 50:
        raise HTTPException(status_code=400, detail="내용이 너무 짧습니다.")

    lecture = Lecture(title=title, content=content)
    db.add(lecture)
    await db.flush()

    # 그래프 먼저 생성 → 노드 목록을 퀴즈 생성에 전달
    nodes = await build_graph(lecture.id, content, db)
    quizzes = await generate_quizzes(lecture.id, content, nodes, db)  # ← nodes 전달

    await db.refresh(lecture)
    return {
        "lecture_id": lecture.id,
        "title": lecture.title,
        "node_count": len(nodes),
        "quiz_count": len(quizzes),
    }


@app.get("/api/lectures")
async def list_lectures(db: AsyncSession = Depends(get_db)):
    result = await db.execute(select(Lecture).order_by(Lecture.created_at.desc()))
    return [
        {"id": l.id, "title": l.title, "created_at": l.created_at.isoformat()}
        for l in result.scalars().all()
    ]


@app.delete("/api/lectures/{lecture_id}")
async def delete_lecture(lecture_id: int, db: AsyncSession = Depends(get_db)):
    lecture = await db.get(Lecture, lecture_id)
    if not lecture:
        raise HTTPException(status_code=404, detail="강의를 찾을 수 없습니다.")
    await db.delete(lecture)
    await db.commit()
    return {"message": "삭제 완료"}


@app.get("/api/lectures/{lecture_id}/graph")
async def get_graph(lecture_id: int, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(GraphNode)
        .where(GraphNode.lecture_id == lecture_id)
        .order_by(GraphNode.level, GraphNode.id)
    )
    nodes = result.scalars().all()
    if not nodes:
        raise HTTPException(status_code=404, detail="그래프 데이터 없음")
    return _build_tree(nodes)


@app.get("/api/lectures/{lecture_id}/quizzes")
async def get_quizzes(lecture_id: int, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(Quiz).where(Quiz.lecture_id == lecture_id).order_by(Quiz.difficulty)
    )
    return [
        {
            "id": q.id,
            "node_id": q.node_id,
            "question": q.question,
            "model_answer": q.model_answer,
            "key_points": json.loads(q.key_points),
            "difficulty": q.difficulty,
            "answer_type": q.answer_type,
            "recommended_time": q.recommended_time,
        }
        for q in result.scalars().all()
    ]


@app.post("/api/quizzes/{quiz_id}/submit")
async def submit_quiz_answer(
    quiz_id: int,
    user_answer: str = Form(...),
    time_spent: int = Form(0),
    db: AsyncSession = Depends(get_db),
):
    quiz = await db.get(Quiz, quiz_id)
    if not quiz:
        raise HTTPException(status_code=404, detail="퀴즈를 찾을 수 없습니다.")

    key_points = json.loads(quiz.key_points)

    # 1. 채점
    try:
        if quiz.answer_type == "quantitative":
            eval_result = evaluate_quantitative_answer(user_answer, quiz.model_answer)
        else:
            eval_result = await evaluate_text_answer(user_answer, quiz.model_answer, key_points)
    except Exception:
        eval_result = {
            "cosine_similarity": 0.0, "keyword_score": 0.0, "final_score": 0.0,
            "matched_keywords": 0, "total_keywords": len(key_points),
            "status": "error", "color": "red", "is_slip": False,
        }

    # 2. 시도 기록
    attempt = QuizAttempt(
        quiz_id=quiz_id,
        lecture_id=quiz.lecture_id,
        user_answer=user_answer,
        cosine_similarity=eval_result.get("cosine_similarity", 0.0),
        keyword_score=eval_result.get("keyword_score", 0.0),
        final_score=eval_result["final_score"],
        status=eval_result["status"],
        is_slip=eval_result.get("is_slip", False),
        time_spent=time_spent,
    )
    db.add(attempt)
    await db.flush()

    llm_analysis = None
    causal_analysis = None

    # 3. 오답 심층 분석
    if eval_result["status"] != "correct":
        try:
            llm_analysis = await multi_agent_analysis(quiz.question, user_answer, quiz.model_answer)
            attempt.error_type = llm_analysis.get("error_type")
        except Exception:
            llm_analysis = {
                "agent_a": "분석 불가", "agent_b": "분석 불가",
                "error_type": "mistake",
                "feedback": "모범 답안과 직접 비교해보세요.",
            }

        try:
            # 퀴즈에 연결된 노드 사용 (없으면 첫 번째 노드 fallback)
            if quiz.node_id:
                target_node = await db.get(GraphNode, quiz.node_id)
            else:
                result = await db.execute(
                    select(GraphNode).where(GraphNode.lecture_id == quiz.lecture_id).limit(1)
                )
                target_node = result.scalar_one_or_none()

            if target_node:
                cos_sim = eval_result.get("cosine_similarity", 0.0)
                causal_analysis = await apply_causal_rules(
                    target_node, eval_result["final_score"], cos_sim, db
                )
                await update_node_weight(
                    target_node, eval_result["final_score"],
                    time_spent, eval_result.get("is_slip", False), db
                )
        except Exception:
            causal_analysis = None

    await db.commit()
    return {
        "attempt_id": attempt.id,
        "evaluation": eval_result,
        "llm_analysis": llm_analysis,
        "causal_analysis": causal_analysis,
    }


@app.get("/api/lectures/{lecture_id}/weak-nodes")
async def get_weak_nodes(lecture_id: int, db: AsyncSession = Depends(get_db)):
    result = await db.execute(
        select(NodeScore, GraphNode)
        .join(GraphNode, NodeScore.node_id == GraphNode.id)
        .where(NodeScore.lecture_id == lecture_id)
        .order_by(NodeScore.final_weight.desc())
    )
    return [
        {
            "node_id": row.GraphNode.id,
            "concept": row.GraphNode.concept,
            "description": row.GraphNode.description,
            "avg_score": round(row.NodeScore.avg_score, 3),
            "final_weight": round(row.NodeScore.final_weight, 3),
            "attempt_count": row.NodeScore.attempt_count,
            "level": row.GraphNode.level,
        }
        for row in result.all()
    ]


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=False)
